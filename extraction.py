# string extraction
from pptx import Presentation
import pandas as pd
import re
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE

# extract text between Double quotes
def getUiString(shape_txt):
    str = re.findall(r'“(.*?)”', shape_txt)
    str += re.findall(r'"(.*?)"', shape_txt)
    return str

# extract text between Double quotes
def getGroupText(groupShape, strings):
    for shape in groupShape.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            getGroupText(shape, strings)
        else:
            if hasattr(shape, "text"):
                strings += getUiString(shape.text)
                # str = re.findall(r'“(.*?)”', shape.text)
                # strings += str
                # strings += re.findall(r'"(.*?)"', shape.text)
    return strings


def extraction(filename):
    prs = Presentation(filename)
    result = []
    i = 0
    for slide in prs.slides:
        i += 1
        if len(slide.shapes) == 0:
            continue
        else:
            title = slide.shapes.title.text

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for r in table.rows:
                        for c in r.cells:
                            txt = ''
                            for t in c.text_frame.paragraphs:
                                txt += t.text
                            strings = getUiString(txt)
                            # strings = re.findall(r'“(.*?)”', txt)
                            # strings += re.findall(r'"(.*?)"', txt)
                            for s in strings:
                                result.append([i, title, s])

                if hasattr(shape, "text"):
                    strings = getUiString(shape.text)
                    # strings = re.findall(r'“(.*?)”', shape.text)
                    # strings += re.findall(r'"(.*?)"', shape.text)
                    #                     print(strings)
                    for s in strings:
                        result.append([i, title, s])

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    strings = []
                    strings = getGroupText(shape, strings)
                    for s in strings:
                        result.append([i, title, s])

    df = pd.DataFrame(result)
    newFileName = filename[:-4] + 'csv'
    os.remove(filename)
    try:
        df.to_csv(newFileName, encoding="euc-kr")
    except:
        df.to_csv(newFileName, encoding="UTF8")
    return newFileName